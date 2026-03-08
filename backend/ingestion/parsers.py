"""
backend/ingestion/parsers.py
─────────────────────────────
Text extraction layer for rich document formats.

Each parser accepts raw bytes and returns a plain UTF-8 string.
The text is then passed to the standard chunker pipeline.

WHY SEPARATE FROM CHUNKER?
  The chunker's job is to *split* text. The parser's job is to *extract* text.
  Separating them makes each module easier to test, replace, and extend.
"""

import io


def extract_pdf(content: bytes) -> str:
    """
    Extract text from a PDF file using pdfplumber.
    pdfplumber handles multi-column layouts and embedded tables better than pypdf2.
    """
    try:
        import pdfplumber
    except ImportError:
        raise RuntimeError("pdfplumber is required. Run: pip install pdfplumber")

    text_parts = []
    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text_parts.append(f"[Page {i + 1}]\n{page_text.strip()}")

    if not text_parts:
        raise ValueError("Could not extract any text from the PDF. It may be image-based.")

    return "\n\n".join(text_parts)


def extract_docx(content: bytes) -> str:
    """
    Extract text from a Word (.docx) file.
    Preserves paragraph structure and table cell content.
    """
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx is required. Run: pip install python-docx")

    doc = Document(io.BytesIO(content))
    text_parts = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # Paragraph text
            text = "".join(node.text or "" for node in element.iter() if node.text)
            if text.strip():
                text_parts.append(text.strip())

        elif tag == "tbl":
            # Table content — join cells row by row
            for row in element:
                row_tag = row.tag.split("}")[-1] if "}" in row.tag else row.tag
                if row_tag == "tr":
                    cell_texts = []
                    for cell in row:
                        cell_tag = cell.tag.split("}")[-1] if "}" in cell.tag else cell.tag
                        if cell_tag == "tc":
                            cell_text = "".join(n.text or "" for n in cell.iter() if n.text)
                            cell_texts.append(cell_text.strip())
                    if cell_texts:
                        text_parts.append(" | ".join(cell_texts))

    if not text_parts:
        raise ValueError("Could not extract any text from the DOCX file.")

    return "\n\n".join(text_parts)


def extract_pptx(content: bytes) -> str:
    """
    Extract text from a PowerPoint (.pptx) file.
    Includes slide text, shapes, and speaker notes.
    """
    try:
        from pptx import Presentation
    except ImportError:
        raise RuntimeError("python-pptx is required. Run: pip install python-pptx")

    prs = Presentation(io.BytesIO(content))
    text_parts = []

    for i, slide in enumerate(prs.slides):
        slide_texts = []

        # Extract text from all shapes (text boxes, titles, etc.)
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_texts.append(shape.text.strip())

        # Extract speaker notes
        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                slide_texts.append(f"[Speaker Notes]: {notes_text}")

        if slide_texts:
            text_parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))

    if not text_parts:
        raise ValueError("Could not extract any text from the PPTX file.")

    return "\n\n".join(text_parts)


# ── Dispatch Table ─────────────────────────────────────────────────────────────

PARSERS: dict[str, callable] = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".pptx": extract_pptx,
}

"""
Supported extensions by category:
  RICH_FORMATS  — requires parsing before chunking (uses this module)
  TEXT_FORMATS  — already plain text, chunked directly
  CODE_FORMATS  — source code, uses structural code chunker
"""
RICH_FORMATS = set(PARSERS.keys())
TEXT_FORMATS = {".md", ".txt", ".rst"}
CODE_FORMATS = {
    # Web / Scripting
    ".py", ".js", ".ts", ".tsx", ".jsx", ".rb", ".php",
    # JVM / Compiled
    ".java", ".kt", ".scala", ".cs",
    # Systems
    ".go", ".rs", ".c", ".cpp", ".cc", ".h", ".hpp",
    # Mobile / Other
    ".swift", ".r", ".lua", ".sh", ".bash",
}
ALL_SUPPORTED = RICH_FORMATS | TEXT_FORMATS | CODE_FORMATS


def extract_text(content: bytes, file_ext: str) -> str:
    """
    Main dispatch function. Given raw file bytes and an extension,
    returns the extracted plain text.
    """
    if file_ext not in RICH_FORMATS:
        raise ValueError(f"No rich parser available for '{file_ext}'.")

    parser_fn = PARSERS[file_ext]
    return parser_fn(content)
